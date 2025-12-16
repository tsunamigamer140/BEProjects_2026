from django.shortcuts import render, redirect, get_object_or_404
from django.contrib.auth import authenticate, login, logout
from django.contrib.auth.decorators import login_required
from django.contrib.admin.views.decorators import staff_member_required
from django.contrib import messages
from django.http import JsonResponse, HttpResponse
from django.views.decorators.csrf import csrf_exempt
from django.contrib.auth.models import User
from django.template.loader import render_to_string
from django.conf import settings
from django.contrib.sites.shortcuts import get_current_site
from django.utils import timezone

from .forms import RegisterForm
from .models import Chat, UserProfile, GuestSession, Subscription
from .rag_llm import rag_answer

from django.core.mail import EmailMultiAlternatives
from email.mime.image import MIMEImage

import os
import uuid
import json
from io import BytesIO

# Optional Export Libraries
try:
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import letter
    REPORTLAB_AVAILABLE = True
except ImportError:
    REPORTLAB_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


# ==============================
# AUTH: Login / Register / Logout
# ==============================

def guest_login(request):
    guest, created = User.objects.get_or_create(
        username="guest_user",
        defaults={"first_name": "Guest"}
    )
    login(request, guest)
    request.session["is_guest"] = True
    return redirect("ai_app:chat")


def login_view(request):
    if request.method == "POST":

        if "guest" in request.POST:
            sid = str(uuid.uuid4())
            GuestSession.objects.create(session_id=sid)
            request.session["guest_id"] = sid
            messages.success(request, "Logged in as Guest!")
            return redirect("ai_app:chat")

        username = request.POST.get("username")
        password = request.POST.get("password")

        user = authenticate(username=username, password=password)
        if user:
            login(request, user)
            UserProfile.objects.get_or_create(user=user)
            messages.success(request, "Login successful!")
            return redirect("ai_app:chat")

        messages.error(request, "Invalid username or password.")

    return render(request, "ai_app/login.html")


def register_view(request):
    just_registered = False

    if request.method == "POST":
        form = RegisterForm(request.POST)
        if form.is_valid():
            user = form.save(commit=False)
            user.set_password(form.cleaned_data["password"])
            user.save()

            UserProfile.objects.create(user=user)

            html_content = render_to_string("ai_app/welcome_email.html", {
                "username": user.username,
                "logo_cid": "roadmaplogo",
            })

            email = EmailMultiAlternatives(
                "Welcome to Roadmap AI & NotebookLM Mini",
                "",
                from_email=settings.EMAIL_HOST_USER,
                to=[user.email]
            )
            email.attach_alternative(html_content, "text/html")

            logo_path = os.path.join(settings.BASE_DIR, "static", "logo.jpg")
            try:
                with open(logo_path, "rb") as f:
                    mime = MIMEImage(f.read())
                    mime.add_header("Content-ID", "<roadmaplogo>")
                    mime.add_header("Content-Disposition", "inline")
                    email.attach(mime)

                email.send()
            except Exception as e:
                print("EMAIL ERROR:", e)

            just_registered = True
            form = RegisterForm()

    else:
        form = RegisterForm()

    return render(request, "ai_app/register.html", {
        "form": form,
        "just_registered": just_registered,
    })


def logout_view(request):
    logout(request)
    return redirect("login")


# ==============================
# CHAT: Queryset Helper
# ==============================

def _get_chat_queryset(request):
    if request.user.is_authenticated:
        return Chat.objects.filter(user=request.user).order_by("timestamp")

    sid = request.session.get("guest_id")
    if not sid:
        return Chat.objects.none()

    return Chat.objects.filter(session_id=sid).order_by("timestamp")


# ==============================
# CHAT APIs
# ==============================

@csrf_exempt
def chat_api(request):
    if request.method != "POST":
        return JsonResponse({"error": "POST only"}, status=400)

    body = json.loads(request.body.decode("utf-8"))
    msg = body.get("message")

    if not msg:
        return JsonResponse({"error": "No message"}, status=400)

    if request.user.is_authenticated:
        profile = UserProfile.objects.get(user=request.user)
        if profile.request_count >= profile.max_requests:
            return JsonResponse({"error": "Limit reached"}, status=403)
        profile.request_count += 1
        profile.save()
        user_or_none = request.user
        session_id = None
    else:
        sid = request.session.get("guest_id")
        guest = GuestSession.objects.get(session_id=sid)
        if guest.request_count >= guest.max_requests:
            return JsonResponse({"error": "Guest limit reached"}, status=403)
        guest.request_count += 1
        guest.save()
        user_or_none = None
        session_id = sid

    response = rag_answer(msg)

    Chat.objects.create(
        user=user_or_none,
        session_id=session_id,
        message=msg,
        response=response
    )

    return JsonResponse({"response": response})


def chat_page(request):
    chats = _get_chat_queryset(request)
    return render(request, "ai_app/chat.html", {"chats": chats})


# Modern RAG API
@csrf_exempt
def rag_chat_api(request):
    if request.method != "POST":
        return JsonResponse({"error": "POST only"}, status=400)

    data = json.loads(request.body.decode("utf-8"))
    query = data.get("query")

    if not query:
        return JsonResponse({"error": "Empty message"}, status=400)

    if request.user.is_authenticated:
        profile = UserProfile.objects.get(user=request.user)
        profile.request_count += 1
        profile.save()
        user_or_none = request.user
        session_id = None
    else:
        sid = request.session.get("guest_id")
        guest = GuestSession.objects.get(session_id=sid)
        guest.request_count += 1
        guest.save()
        user_or_none = None
        session_id = sid

    answer = rag_answer(query)

    Chat.objects.create(
        user=user_or_none,
        session_id=session_id,
        message=query,
        response=answer
    )

    return JsonResponse({"answer": answer})


# ==============================
# CHAT: New / Clear
# ==============================

@csrf_exempt
def new_chat(request):
    if request.method != "POST":
        return JsonResponse({"error": "POST only"}, status=400)

    deleted = _get_chat_queryset(request).delete()[0]
    return JsonResponse({"status": "ok", "deleted": deleted})


@csrf_exempt
def clear_chat(request):
    if request.method != "POST":
        return JsonResponse({"error": "POST only"}, status=400)

    deleted = _get_chat_queryset(request).delete()[0]
    return JsonResponse({"status": "ok", "deleted": deleted})


# ==============================
# PROFILE
# ==============================

@login_required
def profile_page(request):
    profile = UserProfile.objects.get(user=request.user)
    subscription = Subscription.objects.filter(user=request.user).first()   # FIX

    return render(
        request,
        "ai_app/profile.html",
        {
            "profile": profile,
            "subscription": subscription,   # FIX
        }
    )



# ==============================
# SUBSCRIPTION PAGE (User)
# ==============================

@login_required
def subscription_page(request):
    profile = UserProfile.objects.get(user=request.user)
    subscription, _ = Subscription.objects.get_or_create(user=request.user)

    if request.method == "POST" and not subscription.screenshot:
        proof = request.FILES.get("proof")

        if proof:
            subscription.screenshot = proof
            subscription.active = False
            subscription.save()

            dashboard_url = f"http://{get_current_site(request).domain}/app/admin/subscriptions/"
            logo_path = os.path.join(settings.BASE_DIR, "static", "logo.jpg")

            # Admin Email
            admin_html = render_to_string("ai_app/admin_subscription_email.html", {
                "username": request.user.username,
                "user_email": request.user.email,
                "dashboard_url": dashboard_url,
                "logo_cid": "roadmaplogo",
            })

            admin_email = EmailMultiAlternatives(
                "⏳ New Premium Request - Approval Needed",
                "",
                from_email=settings.EMAIL_HOST_USER,
                to=["roadmapai25@gmail.com"],
            )
            admin_email.attach_alternative(admin_html, "text/html")
            admin_email.attach(proof.name, proof.read(), proof.content_type)

            with open(logo_path, "rb") as f:
                mime = MIMEImage(f.read())
                mime.add_header("Content-ID", "<roadmaplogo>")
                admin_email.attach(mime)

            admin_email.send()

            # User Email
            user_html = render_to_string("ai_app/user_subscription_email.html", {
                "username": request.user.username,
                "logo_cid": "roadmaplogo",
            })

            user_email = EmailMultiAlternatives(
                "✅ Payment Submitted — Pending Verification",
                "",
                from_email=settings.EMAIL_HOST_USER,
                to=[request.user.email],
            )
            user_email.attach_alternative(user_html, "text/html")

            with open(logo_path, "rb") as f:
                mime2 = MIMEImage(f.read())
                mime2.add_header("Content-ID", "<roadmaplogo>")
                user_email.attach(mime2)

            user_email.send()

            messages.success(request, "Your payment was submitted. You will receive an update soon.")
            return redirect("ai_app:subscription")

    already_submitted = bool(subscription.screenshot)

    return render(request, "ai_app/subscription.html", {
        "profile": profile,
        "subscription": subscription,
        "already_submitted": already_submitted,
    })


# ==============================
# ADMIN DASHBOARD: Subscription List
# ==============================

@staff_member_required
def admin_subscriptions(request):
    subs = Subscription.objects.all().order_by("-requested_at")
    return render(request, "ai_app/admin_subscriptions.html", {"subs": subs})


# ==============================
# ADMIN: Approve Subscription
# ==============================

@staff_member_required
def approve_subscription(request, pk):
    sub = get_object_or_404(Subscription, id=pk)
    profile = UserProfile.objects.get(user=sub.user)

    sub.active = True
    sub.approved_at = timezone.now()
    sub.save()

    profile.max_requests = 999999
    profile.save()

    html = render_to_string("ai_app/subscription_approved_email.html", {
        "username": sub.user.username,
        "logo_cid": "roadmaplogo",
    })

    email = EmailMultiAlternatives(
        "🎉 Premium Subscription Approved!",
        "",
        from_email=settings.EMAIL_HOST_USER,
        to=[sub.user.email],
    )
    email.attach_alternative(html, "text/html")

    logo_path = os.path.join(settings.BASE_DIR, "static", "logo.jpg")
    with open(logo_path, "rb") as f:
        mime = MIMEImage(f.read())
        mime.add_header("Content-ID", "<roadmaplogo>")
        email.attach(mime)

    email.send()

    messages.success(request, "Subscription approved successfully!")
    return redirect("ai_app:admin_subscriptions")


# ==============================
# ADMIN: Reject Subscription
# ==============================

@staff_member_required
def reject_subscription(request, pk):
    sub = get_object_or_404(Subscription, id=pk)
    sub.active = False
    sub.screenshot = None  # this marks it as REJECTED
    sub.save()

    html = render_to_string("ai_app/subscription_rejected_email.html", {
        "username": sub.user.username,
        "logo_cid": "roadmaplogo",
    })

    email = EmailMultiAlternatives(
        "❌ Premium Subscription Rejected",
        "",
        from_email=settings.EMAIL_HOST_USER,
        to=[sub.user.email],
    )
    email.attach_alternative(html, "text/html")

    logo_path = os.path.join(settings.BASE_DIR, "static", "logo.jpg")
    with open(logo_path, "rb") as f:
        mime = MIMEImage(f.read())
        mime.add_header("Content-ID", "<roadmaplogo>")
        email.attach(mime)

    email.send()

    messages.error(request, "Subscription rejected.")
    return redirect("ai_app:admin_subscriptions")


# ==============================
# EXPORT CHAT
# ==============================

def export_chat(request, fmt):
    qs = _get_chat_queryset(request)
    lines = []

    for c in qs:
        lines.append(f"You: {c.message}")
        lines.append(f"AI: {c.response}")
        lines.append("")

    text_content = "\n".join(lines)

    if fmt == "txt":
        resp = HttpResponse(text_content, content_type="text/plain; charset=utf-8")
        resp["Content-Disposition"] = 'attachment; filename="chat_history.txt"'
        return resp

    if fmt == "json":
        data = [{
            "message": c.message,
            "response": c.response,
            "timestamp": c.timestamp.isoformat()
        } for c in qs]

        resp = HttpResponse(json.dumps(data, indent=2), content_type="application/json")
        resp["Content-Disposition"] = 'attachment; filename="chat_history.json"'
        return resp

    if fmt == "pdf":
        if not REPORTLAB_AVAILABLE:
            return HttpResponse("reportlab not installed.", status=500)

        buffer = BytesIO()
        p = canvas.Canvas(buffer, pagesize=letter)
        w, h = letter
        y = h - 50

        for line in text_content.split("\n"):
            p.drawString(40, y, line)
            y -= 15
            if y < 40:
                p.showPage()
                y = h - 50

        p.save()
        pdf = buffer.getvalue()
        buffer.close()

        resp = HttpResponse(pdf, content_type="application/pdf")
        resp["Content-Disposition"] = 'attachment; filename="chat_history.pdf"'
        return resp

    if fmt == "pptx":
        if not PPTX_AVAILABLE:
            return HttpResponse("python-pptx not installed.", status=500)

        prs = Presentation()
        layout = prs.slide_layouts[1]

        for chat in qs:
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = f"You: {chat.message[:40]}..."
            slide.placeholders[1].text = f"You:\n{chat.message}\n\nAI:\n{chat.response}"

        stream = BytesIO()
        prs.save(stream)
        stream.seek(0)

        resp = HttpResponse(
            stream.read(),
            content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
        resp["Content-Disposition"] = 'attachment; filename="chat_history.pptx"'
        return resp

    return HttpResponse("Unsupported format", status=400)
